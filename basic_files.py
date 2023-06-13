import random
import os
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from faker import Faker

fake = Faker()

numberOfFiles = 10

# Define different parts of the filename
base_name = ("project", "new","file","workplace", "report", "training", "meeting" ,"screenshot", "09-28-2021")
extra_name = ("Q1","Q2", "Q3", "Q4", "<COMPANY_NAME>", "personal","school")
file_ext = (".txt", ".pdf", ".docx")
all_names = (extra_name, base_name, file_ext)


for n in range(numberOfFiles):
    ext = random.choice(file_ext)
    # Generate a random filename by choosing parts from different tuples
    filename = ' '.join([random.choice(i) for i in all_names])
    
    if ext == '.pdf':
        # Generate a PDF file
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        # Add random numbers to the PDF
        for _ in range(10):
            pdf.cell(0, 10, txt=str(random.randint(0,60 + (n * 10))), ln=True)
        
        # Add a placeholder message to the PDF
        pdf.cell(0, 10, txt="ADD COMPANY NAME AND/OR CUSTOM MESSAGE HERE.", ln=True)
        
        # Save the PDF with the generated filename
        pdf.output(filename)
    
    elif ext == '.docx':
        # Generate a Word document
        doc = Document()
        
        # Add random numbers as paragraphs to the document
        for _ in range(10):
            doc.add_paragraph(str(random.randint(0,60 + (n * 10))))
        
        # Add a placeholder message as a paragraph to the document
        doc.add_paragraph("ADD COMPANY NAME AND/OR CUSTOM MESSAGE HERE.")
        
        # Save the document with the generated filename
        doc.save(filename)
    
    elif ext == '.pptx':
        # Generate a PowerPoint presentation
        pres = Presentation()
        slide_layout = pres.slide_layouts[0]  # choosing a slide layout
        slide = pres.slides.add_slide(slide_layout)  # adding a slide
        title = slide.shapes.title  # defining title
        subtitle = slide.placeholders[1]  # defining subtitle
        
        # Add random numbers to slides
        for _ in range(10):
            subtitle.text = str(random.randint(0,60 + (n * 10)))
            slide = pres.slides.add_slide(slide_layout)
            subtitle = slide.placeholders[1]
        
        # Add a placeholder message to the final slide
        subtitle.text = "ADD COMPANY NAME AND/OR CUSTOM MESSAGE HERE."
        
        # Save the presentation with the generated filename
        pres.save(filename)
    
    elif ext == '.xlsx':
        # Generate an Excel workbook
        workbook = Workbook()
        sheet = workbook.active

        # Write headers to the worksheet
        headers = ["First Name", "Last Name", "Address", "Phone Number"]
        sheet.append(headers)

        # Generate fake data and write to the worksheet
        for _ in range(10):
            first_name = fake.first_name()
            last_name = fake.last_name()
            address = fake.address().replace('\n', ', ')
            phone_number = fake.phone_number()
            row_data = [first_name, last_name, address, phone_number]
            sheet.append(row_data)

        # Save the workbook with the generated filename
        workbook.save(filename)
