import random
import os
from fpdf import FPDF
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from faker import Faker

fake = Faker()

numberOfFiles = 10


base_name = ("project", "new","file","workplace", "report", "training", "meeting" ,"screenshot", "09-28-2021")
extra_name = ("Q2", "Q3", "Q4", "VMC", "NSX", "WorkSpace One", "personal","vmware","school")
file_ext = (".txt", ".pdf", ".docx", ".xlsx")
all_names = (extra_name, base_name, file_ext)


for n in range(numberOfFiles):
    ext = random.choice(file_ext)
    filename = ' '.join([random.choice(i) for i in all_names])
    
    if ext == '.pdf':
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for _ in range(10):
            pdf.cell(0, 10, txt=str(random.randint(0,60 + (n * 10))), ln=True)
        pdf.cell(0, 10, txt="Copyright (c) 1998-2009 VMware, Inc.  All rights reserved.", ln=True)
        pdf.output(filename)
    elif ext == '.docx':
        doc = Document()
        for _ in range(10):
            doc.add_paragraph(str(random.randint(0,60 + (n * 10))))
        doc.add_paragraph("Copyright (c) 1998-2009 VMware, Inc.  All rights reserved.")
        doc.save(filename)
    elif ext == '.pptx':
        pres = Presentation()
        slide_layout = pres.slide_layouts[0]  # choosing a slide layout
        slide = pres.slides.add_slide(slide_layout)  # adding a slide
        title = slide.shapes.title  # defining title
        subtitle = slide.placeholders[1]  # defining subtitle
        
        for _ in range(10):
            subtitle.text = str(random.randint(0,60 + (n * 10)))
            slide = pres.slides.add_slide(slide_layout)
            subtitle = slide.placeholders[1]
        
        subtitle.text = "Copyright (c) 1998-2009 VMware, Inc.  All rights reserved."
        pres.save(filename)
    elif ext == '.xlsx':
        workbook = Workbook()
        sheet = workbook.active

        # write headers
        headers = ["First Name", "Last Name", "Address", "Phone Number"]
        sheet.append(headers)

        # write fake data
        for _ in range(10):
            first_name = fake.first_name()
            last_name = fake.last_name()
            address = fake.address().replace('\n', ', ')
            phone_number = fake.phone_number()
            row_data = [first_name, last_name, address, phone_number]
            sheet.append(row_data)

        workbook.save(filename)
    else:
        with open(filename, 'w') as newFile:
            for _ in range(10):
                newFile.write(str(random.randint(0,60 + (n * 10))) + "\n")
            newFile.write("Copyright (c) 1998-2009 VMware, Inc.  All rights reserved.\n")
